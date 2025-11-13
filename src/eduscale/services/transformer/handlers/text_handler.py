"""Text extraction handler for various document formats."""

import logging
import subprocess
from datetime import datetime
from pathlib import Path
from typing import NamedTuple

import yaml
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
from odf import text as odf_text, teletype
from odf.opendocument import load as load_odf
from striprtf.striprtf import rtf_to_text

from eduscale.services.transformer.exceptions import ExtractionError

logger = logging.getLogger(__name__)


class ExtractionMetadata(NamedTuple):
    """Metadata about the extraction process."""
    extraction_method: str
    page_count: int | None = None
    sheet_count: int | None = None
    slide_count: int | None = None
    word_count: int | None = None
    character_count: int | None = None


def build_text_frontmatter(
    file_id: str,
    region_id: str,
    text_uri: str,
    file_category: str,
    extraction_metadata: ExtractionMetadata,
    original_filename: str | None = None,
    original_content_type: str | None = None,
    original_size_bytes: int | None = None,
    bucket: str | None = None,
    object_path: str | None = None,
    event_id: str | None = None,
    uploaded_at: str | None = None,
    extraction_duration_ms: int | None = None,
) -> str:
    """Build YAML frontmatter with metadata for AI processing.

    Args:
        file_id: Unique file identifier
        region_id: Region identifier
        text_uri: GCS URI of the extracted text file
        file_category: File category (text, audio, etc.)
        extraction_metadata: Metadata from the extraction process
        original_filename: Original filename
        original_content_type: Original MIME type
        original_size_bytes: Original file size in bytes
        bucket: GCS bucket name
        object_path: Full object path in GCS
        event_id: CloudEvent ID for tracing
        uploaded_at: Upload timestamp (ISO format)
        extraction_duration_ms: Time taken for extraction in milliseconds

    Returns:
        YAML frontmatter string with metadata
    """
    # Build metadata dictionary
    metadata = {
        "file_id": file_id,
        "region_id": region_id,
        "text_uri": text_uri,
    }

    # Add event ID if available
    if event_id:
        metadata["event_id"] = event_id

    # Add original file information
    original = {}
    if original_filename:
        original["filename"] = original_filename
    if original_content_type:
        original["content_type"] = original_content_type
    if original_size_bytes is not None:
        original["size_bytes"] = original_size_bytes
    if bucket:
        original["bucket"] = bucket
    if object_path:
        original["object_path"] = object_path
    if uploaded_at:
        original["uploaded_at"] = uploaded_at

    if original:
        metadata["original"] = original

    # Add file category
    metadata["file_category"] = file_category

    # Add extraction information
    extraction = {
        "method": extraction_metadata.extraction_method,
        "timestamp": datetime.utcnow().isoformat() + "Z",
        "success": True,
    }
    if extraction_duration_ms is not None:
        extraction["duration_ms"] = extraction_duration_ms

    metadata["extraction"] = extraction

    # Add content metrics
    content = {}
    if extraction_metadata.character_count is not None:
        content["text_length"] = extraction_metadata.character_count
    if extraction_metadata.word_count is not None:
        content["word_count"] = extraction_metadata.word_count
    if extraction_metadata.character_count is not None:
        content["character_count"] = extraction_metadata.character_count

    if content:
        metadata["content"] = content

    # Add document-specific metadata
    document = {}
    if extraction_metadata.page_count is not None:
        document["page_count"] = extraction_metadata.page_count
    if extraction_metadata.sheet_count is not None:
        document["sheet_count"] = extraction_metadata.sheet_count
    if extraction_metadata.slide_count is not None:
        document["slide_count"] = extraction_metadata.slide_count

    if document:
        metadata["document"] = document

    # Convert to YAML with proper formatting
    yaml_content = yaml.dump(
        metadata,
        default_flow_style=False,
        allow_unicode=True,
        sort_keys=False,
    )

    # Return with frontmatter delimiters
    return f"---\n{yaml_content}---\n"


def extract_text_from_pdf(file_path: Path) -> tuple[str, ExtractionMetadata]:
    """Extract text from PDF using pdfplumber.

    Args:
        file_path: Path to the PDF file

    Returns:
        Tuple of (extracted text, metadata)

    Raises:
        ExtractionError: If extraction fails
    """
    try:
        logger.info("Extracting text from PDF", extra={"file_path": str(file_path)})

        text_parts = []
        page_count = 0

        with pdfplumber.open(file_path) as pdf:
            page_count = len(pdf.pages)
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)

        text = "\n\n".join(text_parts)
        word_count = len(text.split())

        metadata = ExtractionMetadata(
            extraction_method="pdfplumber",
            page_count=page_count,
            word_count=word_count,
            character_count=len(text),
        )

        logger.info(
            "PDF extraction successful",
            extra={
                "file_path": str(file_path),
                "page_count": page_count,
                "word_count": word_count,
            },
        )

        return text, metadata
    except Exception as e:
        logger.error(
            "Failed to extract text from PDF",
            extra={"file_path": str(file_path), "error": str(e)},
        )
        raise ExtractionError(f"PDF extraction failed: {e}") from e


def extract_text_from_docx(file_path: Path) -> tuple[str, ExtractionMetadata]:
    """Extract text from DOCX using python-docx.

    Args:
        file_path: Path to the DOCX file

    Returns:
        Tuple of (extracted text, metadata)

    Raises:
        ExtractionError: If extraction fails
    """
    try:
        logger.info("Extracting text from DOCX", extra={"file_path": str(file_path)})

        doc = DocxDocument(file_path)
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        text = "\n\n".join(paragraphs)

        word_count = len(text.split())

        metadata = ExtractionMetadata(
            extraction_method="python-docx",
            word_count=word_count,
            character_count=len(text),
        )

        logger.info(
            "DOCX extraction successful",
            extra={"file_path": str(file_path), "word_count": word_count},
        )

        return text, metadata
    except Exception as e:
        logger.error(
            "Failed to extract text from DOCX",
            extra={"file_path": str(file_path), "error": str(e)},
        )
        raise ExtractionError(f"DOCX extraction failed: {e}") from e


def extract_text_from_xlsx(file_path: Path) -> tuple[str, ExtractionMetadata]:
    """Extract text from XLSX using openpyxl.

    Args:
        file_path: Path to the XLSX file

    Returns:
        Tuple of (extracted text, metadata)

    Raises:
        ExtractionError: If extraction fails
    """
    try:
        logger.info("Extracting text from XLSX", extra={"file_path": str(file_path)})

        workbook = load_workbook(file_path, read_only=True, data_only=True)
        text_parts = []
        sheet_count = len(workbook.sheetnames)

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"=== Sheet: {sheet_name} ===\n")

            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    text_parts.append(row_text)

        text = "\n".join(text_parts)
        word_count = len(text.split())

        metadata = ExtractionMetadata(
            extraction_method="openpyxl",
            sheet_count=sheet_count,
            word_count=word_count,
            character_count=len(text),
        )

        logger.info(
            "XLSX extraction successful",
            extra={
                "file_path": str(file_path),
                "sheet_count": sheet_count,
                "word_count": word_count,
            },
        )

        return text, metadata
    except Exception as e:
        logger.error(
            "Failed to extract text from XLSX",
            extra={"file_path": str(file_path), "error": str(e)},
        )
        raise ExtractionError(f"XLSX extraction failed: {e}") from e


def extract_text_from_pptx(file_path: Path) -> tuple[str, ExtractionMetadata]:
    """Extract text from PPTX using python-pptx.

    Args:
        file_path: Path to the PPTX file

    Returns:
        Tuple of (extracted text, metadata)

    Raises:
        ExtractionError: If extraction fails
    """
    try:
        logger.info("Extracting text from PPTX", extra={"file_path": str(file_path)})

        prs = Presentation(file_path)
        text_parts = []
        slide_count = len(prs.slides)

        for i, slide in enumerate(prs.slides, 1):
            text_parts.append(f"=== Slide {i} ===\n")

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)

            # Extract notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes:
                    text_parts.append(f"\nNotes: {notes}")

        text = "\n\n".join(text_parts)
        word_count = len(text.split())

        metadata = ExtractionMetadata(
            extraction_method="python-pptx",
            slide_count=slide_count,
            word_count=word_count,
            character_count=len(text),
        )

        logger.info(
            "PPTX extraction successful",
            extra={
                "file_path": str(file_path),
                "slide_count": slide_count,
                "word_count": word_count,
            },
        )

        return text, metadata
    except Exception as e:
        logger.error(
            "Failed to extract text from PPTX",
            extra={"file_path": str(file_path), "error": str(e)},
        )
        raise ExtractionError(f"PPTX extraction failed: {e}") from e


def extract_text_from_odt(file_path: Path) -> tuple[str, ExtractionMetadata]:
    """Extract text from ODT (OpenDocument Text) using odfpy.

    Args:
        file_path: Path to the ODT file

    Returns:
        Tuple of (extracted text, metadata)

    Raises:
        ExtractionError: If extraction fails
    """
    try:
        logger.info("Extracting text from ODT", extra={"file_path": str(file_path)})

        doc = load_odf(str(file_path))
        text_parts = []

        for para in doc.getElementsByType(odf_text.P):
            text_parts.append(teletype.extractText(para))

        text = "\n\n".join(text_parts)
        word_count = len(text.split())

        metadata = ExtractionMetadata(
            extraction_method="odfpy",
            word_count=word_count,
            character_count=len(text),
        )

        logger.info(
            "ODT extraction successful",
            extra={"file_path": str(file_path), "word_count": word_count},
        )

        return text, metadata
    except Exception as e:
        logger.error(
            "Failed to extract text from ODT",
            extra={"file_path": str(file_path), "error": str(e)},
        )
        raise ExtractionError(f"ODT extraction failed: {e}") from e


def extract_text_from_ods(file_path: Path) -> tuple[str, ExtractionMetadata]:
    """Extract text from ODS (OpenDocument Spreadsheet) using odfpy.

    Args:
        file_path: Path to the ODS file

    Returns:
        Tuple of (extracted text, metadata)

    Raises:
        ExtractionError: If extraction fails
    """
    try:
        logger.info("Extracting text from ODS", extra={"file_path": str(file_path)})

        doc = load_odf(str(file_path))
        from odf.table import Table, TableRow, TableCell

        text_parts = []
        tables = doc.getElementsByType(Table)
        sheet_count = len(tables)

        for table in tables:
            table_name = table.getAttribute("name") or "Unnamed Sheet"
            text_parts.append(f"=== Sheet: {table_name} ===\n")

            rows = table.getElementsByType(TableRow)
            for row in rows:
                cells = row.getElementsByType(TableCell)
                cell_texts = [teletype.extractText(cell) for cell in cells]
                row_text = "\t".join(cell_texts)
                if row_text.strip():
                    text_parts.append(row_text)

        text = "\n".join(text_parts)
        word_count = len(text.split())

        metadata = ExtractionMetadata(
            extraction_method="odfpy",
            sheet_count=sheet_count,
            word_count=word_count,
            character_count=len(text),
        )

        logger.info(
            "ODS extraction successful",
            extra={
                "file_path": str(file_path),
                "sheet_count": sheet_count,
                "word_count": word_count,
            },
        )

        return text, metadata
    except Exception as e:
        logger.error(
            "Failed to extract text from ODS",
            extra={"file_path": str(file_path), "error": str(e)},
        )
        raise ExtractionError(f"ODS extraction failed: {e}") from e


def extract_text_from_odp(file_path: Path) -> tuple[str, ExtractionMetadata]:
    """Extract text from ODP (OpenDocument Presentation) using odfpy.

    Args:
        file_path: Path to the ODP file

    Returns:
        Tuple of (extracted text, metadata)

    Raises:
        ExtractionError: If extraction fails
    """
    try:
        logger.info("Extracting text from ODP", extra={"file_path": str(file_path)})

        doc = load_odf(str(file_path))
        from odf.draw import Page

        text_parts = []
        pages = doc.getElementsByType(Page)
        slide_count = len(pages)

        for i, page in enumerate(pages, 1):
            text_parts.append(f"=== Slide {i} ===\n")
            text_parts.append(teletype.extractText(page))

        text = "\n\n".join(text_parts)
        word_count = len(text.split())

        metadata = ExtractionMetadata(
            extraction_method="odfpy",
            slide_count=slide_count,
            word_count=word_count,
            character_count=len(text),
        )

        logger.info(
            "ODP extraction successful",
            extra={
                "file_path": str(file_path),
                "slide_count": slide_count,
                "word_count": word_count,
            },
        )

        return text, metadata
    except Exception as e:
        logger.error(
            "Failed to extract text from ODP",
            extra={"file_path": str(file_path), "error": str(e)},
        )
        raise ExtractionError(f"ODP extraction failed: {e}") from e


def extract_text_from_rtf(file_path: Path) -> tuple[str, ExtractionMetadata]:
    """Extract text from RTF using striprtf.

    Args:
        file_path: Path to the RTF file

    Returns:
        Tuple of (extracted text, metadata)

    Raises:
        ExtractionError: If extraction fails
    """
    try:
        logger.info("Extracting text from RTF", extra={"file_path": str(file_path)})

        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            rtf_content = f.read()

        text = rtf_to_text(rtf_content)
        word_count = len(text.split())

        metadata = ExtractionMetadata(
            extraction_method="striprtf",
            word_count=word_count,
            character_count=len(text),
        )

        logger.info(
            "RTF extraction successful",
            extra={"file_path": str(file_path), "word_count": word_count},
        )

        return text, metadata
    except Exception as e:
        logger.error(
            "Failed to extract text from RTF",
            extra={"file_path": str(file_path), "error": str(e)},
        )
        raise ExtractionError(f"RTF extraction failed: {e}") from e


def extract_text_from_doc(file_path: Path) -> tuple[str, ExtractionMetadata]:
    """Extract text from legacy DOC format using antiword.

    Args:
        file_path: Path to the DOC file

    Returns:
        Tuple of (extracted text, metadata)

    Raises:
        ExtractionError: If extraction fails
    """
    try:
        logger.info("Extracting text from DOC", extra={"file_path": str(file_path)})

        # Try using antiword command-line tool
        result = subprocess.run(
            ["antiword", str(file_path)],
            capture_output=True,
            text=True,
            timeout=30,
        )

        if result.returncode != 0:
            raise ExtractionError(f"antiword failed: {result.stderr}")

        text = result.stdout
        word_count = len(text.split())

        metadata = ExtractionMetadata(
            extraction_method="antiword",
            word_count=word_count,
            character_count=len(text),
        )

        logger.info(
            "DOC extraction successful",
            extra={"file_path": str(file_path), "word_count": word_count},
        )

        return text, metadata
    except FileNotFoundError:
        logger.error("antiword not installed")
        raise ExtractionError("antiword not installed. Install with: apt-get install antiword")
    except Exception as e:
        logger.error(
            "Failed to extract text from DOC",
            extra={"file_path": str(file_path), "error": str(e)},
        )
        raise ExtractionError(f"DOC extraction failed: {e}") from e


def extract_text_from_plain(file_path: Path) -> tuple[str, ExtractionMetadata]:
    """Extract text from plain text files with encoding detection.

    Args:
        file_path: Path to the text file

    Returns:
        Tuple of (extracted text, metadata)

    Raises:
        ExtractionError: If extraction fails
    """
    try:
        logger.info("Extracting text from plain file", extra={"file_path": str(file_path)})

        # Try UTF-8 first
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
        except UnicodeDecodeError:
            # Fallback to latin-1
            logger.warning("UTF-8 failed, trying latin-1", extra={"file_path": str(file_path)})
            with open(file_path, "r", encoding="latin-1") as f:
                text = f.read()

        word_count = len(text.split())

        metadata = ExtractionMetadata(
            extraction_method="plain_text",
            word_count=word_count,
            character_count=len(text),
        )

        logger.info(
            "Plain text extraction successful",
            extra={"file_path": str(file_path), "word_count": word_count},
        )

        return text, metadata
    except Exception as e:
        logger.error(
            "Failed to extract text from plain file",
            extra={"file_path": str(file_path), "error": str(e)},
        )
        raise ExtractionError(f"Plain text extraction failed: {e}") from e


def extract_text(file_path: Path, content_type: str) -> tuple[str, ExtractionMetadata]:
    """Route text extraction based on content type.

    Args:
        file_path: Path to the file
        content_type: MIME type of the file

    Returns:
        Tuple of (extracted text, metadata)

    Raises:
        ExtractionError: If extraction fails or format not supported
    """
    logger.info(
        "Routing text extraction",
        extra={"file_path": str(file_path), "content_type": content_type},
    )

    # Map content types to extraction functions
    extractors = {
        "application/pdf": extract_text_from_pdf,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": extract_text_from_docx,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": extract_text_from_xlsx,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": extract_text_from_pptx,
        "application/vnd.oasis.opendocument.text": extract_text_from_odt,
        "application/vnd.oasis.opendocument.spreadsheet": extract_text_from_ods,
        "application/vnd.oasis.opendocument.presentation": extract_text_from_odp,
        "application/rtf": extract_text_from_rtf,
        "application/msword": extract_text_from_doc,
        "text/plain": extract_text_from_plain,
        "text/markdown": extract_text_from_plain,
        "text/html": extract_text_from_plain,
    }

    extractor = extractors.get(content_type)
    if not extractor:
        raise ExtractionError(f"Unsupported content type: {content_type}")

    return extractor(file_path)
